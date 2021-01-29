import subprocess
from selenium import webdriver
from webdriver_manager.chrome import ChromeDriverManager
import os
from bs4 import BeautifulSoup
import time
import getpass
from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO


def func_pcn(l):
    print("[INFO] Shutting down all chrome processes to use this program.")
    subprocess.call("TASKKILL /f  /IM  CHROME.EXE")
    #import getpass
    options = webdriver.ChromeOptions()
    #Find your google cookie directory. I have automated this part, if it doesn't work please refer to my video to find your cookie folder.
    path_to_chrome_cookie="user-data-dir=C:\\Users\\rajdchak\\AppData\\Local\\Google\\Chrome\\User Data"
    #Path to your chrome profile
    options.add_argument(path_to_chrome_cookie) 

    driver = webdriver.Chrome(ChromeDriverManager().install(), options=options)
    #output = BytesIO()
    prs = Presentation()
    output = BytesIO()
    
    #num=['302266','302269','302312','302235','302311']
    title=["Summary :","Reason for Change :","Description of Change : ","Effect of Change : "]

    #go to the website that you want to bypass
    number=0
    for n in l:
        count=0

        number=number+1
        
        url="https://pcncle.cloudapps.cisco.com/emco/pcnclei/prsc/pcnReport.do?recordManagementView="+n
        driver.get(url)
        content = driver.page_source
        soup = BeautifulSoup(content)
        #print(soup)

        blank_slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(blank_slide_layout)
        left = top = width = height = Inches(0.2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text ="PCN "+str(number)
        p.font.bold = True
            
        for a in soup.findAll('strong'):
            print(a.text)
        for b in soup.findAll('div'):
            print(b.text)
            
            """p = tf.add_paragraph()
            p.font.bold = True
            p.text = title[count]"""
            p = tf.add_paragraph()
            p.text = title[count]
            p.font.bold = True
            
            p = tf.add_paragraph()
            p.text = b.text+"\n\n"
        
            
            #tf.text = tf.text+p.text
            #tf.text = tf.text+"\n"+b.text+"\n\n"
            count=count+1
            if(count==2):
                
                p = tf.add_paragraph()
                p.text = "Please refer to the Part Change Notification at the following url: "+"\n"
                
                run = p.add_run()
                run.text = url
                hlink = run.hyperlink
                hlink.address = url
                run.hyperlink.address=url
            if(count==4):
                break
    filename="pcn.pptx"
    prs.save(filename)
    with open('pcn.pptx', 'rb') as f:
        source_stream = BytesIO(f.read())

    ppt = source_stream.getvalue()
    source_stream.close()
    return ppt
    